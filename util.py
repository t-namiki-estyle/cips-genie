import logging
import urllib.request
import urllib.parse
import json
import os
import re
import io
import time
import base64
import random
import asyncio
import ast
import uuid

from copy import deepcopy
from datetime import datetime, timedelta
from zoneinfo import ZoneInfo
from mimetypes import guess_type

import smtplib
from email.message import EmailMessage
from email.header import Header
from email.utils import formataddr
from email.utils import make_msgid

# OSSライブラリ
import requests
from requests.exceptions import RequestException, Timeout, TooManyRedirects
import openpyxl
import pandas as pd
import tiktoken
import openai
from httpx import Timeout
import extract_msg
from bs4 import BeautifulSoup, Comment

# import collections.abc # インポートしないとエラーが発生する
from pptx import Presentation
import docx
from docx.oxml.table import CT_Tbl
from docx.table import Table
from docx.text.paragraph import Paragraph

from azure.storage.blob import BlobServiceClient
from azure.cosmos import CosmosClient, PartitionKey
from azure.cosmos.exceptions import CosmosResourceExistsError
from azure.ai.formrecognizer import DocumentAnalysisClient
from azure.core.credentials import AzureKeyCredential

# 自作
from config import (LLM_REGISTRY, COSMOS_CONNECTION_STRING, BLOB_CONNECTION_STRING,
                    AUDIO_CONTAINER_NAME, FILE_CONTAINER_NAME, VARIABLE_LIST, ENVIRONMENT_SELECTED)
from i_style.aiohttp import AsyncHttpClient
from i_style.llm import AzureOpenAI
from i_style.text_extractor import FileTextExtractor

############
# get path #
############
# Google
GOOGLE_API_KEY = os.environ.get("GOOGLE_API_KEY")
GOOGLE_CSE_ID = os.environ.get("GOOGLE_CSE_ID")

# 文字エンコード
ENCODINGS = ['utf-8', 'shift_jis', 'euc_jp',
             'iso2022_jp', 'cp932', 'utf-16', 'latin1']

##########
# module #
##########

##############
# google api #
##############


async def google_search(async_http_client: AsyncHttpClient, query: str, num: int = 5) -> list[dict]:
    """
    Google検索で上位 num 件を取得し、[{'title': ..., 'snippet': ..., 'link': ...}, ...] を返す関数
    """
    base_url = "https://www.googleapis.com/customsearch/v1"
    # GET パラメータとして API キーやその他オプションを渡す
    params = {
        "q": query,
        "key": GOOGLE_API_KEY,
        "cx": GOOGLE_CSE_ID,
        "num": num,
        "gl": "jp",
        # "dateRestrict": "y1",
    }

    try:
        resp = await async_http_client.get(
            base_url,
            params=params,
            process_name="GoogleCustomSearch",
            raise_for_status=True,
        )
    except Exception as e:
        logging.error(f"Google Search Error: {e}")
        return [{"title": "ERROR",
                "snippet": "Google APIに接続できませんでした。",
                 "link": ""
                 }]

    # レスポンスから items を取得
    items = resp.get("items", [])
    if not items:
        return [{"title": "ERROR",
                "snippet": "Google検索ではお探しの情報を見つけられませんでした。",
                 "link": ""
                 }]

    # 必要なフィールドだけ抜き出して返す
    return [{"title": it.get("title", ""),
            "snippet": it.get("snippet", ""),
             "link": it.get("link", "")}
            for it in items]

# urlから必要な情報のみ抽出する関数


async def url2text(url: str, max_length: int = 20_000, timeout: int = 10) -> str:
    """
    urlを受け取り、htmlを取得、最低限の整形を行なった内容をstrで返却する

    Args:
        url: 取得対象のURL
        max_length: 返却するテキストの最大長
        timeout: リクエストのタイムアウト秒数

    Returns:
        整形されたHTML内容、またはエラーメッセージ
    """
    # URLの検証
    if not is_valid_url(url):
        return "無効なURLフォーマットです。"

    # 内容の取得
    if url.split(".")[-1].lower() == "pdf":
        text = await pdf2text(url, max_length, timeout)
    else:
        text = html2text(url, max_length, timeout)

    return text

# 上記の関数を用いてcse_listを更新する関数


async def updateSearchResults(cse_list: list, max_length: int = 20_000) -> list:
    """
    google/bingの検索結果からより詳細な内容を取得する
    snippetに追記する形で更新し、詳細取得に失敗した場合はsnippetを更新しない

    Args:
        cse_list: 検索結果のリスト

    Returns:
        更新された検索結果のリスト
    """
    if not isinstance(cse_list, list):
        return []

    for item in cse_list:
        try:
            if not isinstance(item, dict) or "link" not in item or "snippet" not in item:
                continue

            # html_content取得
            html_content = await url2text(item["link"], max_length)

            # エラーメッセージが返ってきた場合はスキップ
            if html_content.startswith("URL") or html_content.startswith("無効") or html_content.startswith("予期せぬ"):
                continue

            # 要約 + 本文
            item["snippet"] += "\n\n## 本文\n\n" + html_content
        except Exception as e:
            # 個別のアイテム処理でエラーが発生しても全体の処理は継続
            continue

    return cse_list


async def pdf2text(url: str, max_length: int = 20_000, timeout: int = 10) -> str:
    """
    pdfのurlを受け取り、htmlを取得、不要なタグを削除し整形した内容をstrで返却する

    Args:
        url: 取得対象のURL
        max_length: 返却するテキストの最大長
        timeout: リクエストのタイムアウト秒数

    Returns:
        整形されたHTML内容、またはエラーメッセージ
    """
    # LLM用の環境変数を設定（本番では適切な値を設定）

    # FileTextExtractorの初期化
    extractor = FileTextExtractor(
        file_extension_list=["pdf",],
    )

    headers = {
        'User-Agent': 'Mozilla/5.0 (compatible; MyWebScraper/1.0)',
        'Accept': 'application/pdf',
        'Accept-Language': 'ja,en-US;q=0.9,en;q=0.8',
    }
    try:
        response = requests.get(url, headers=headers,
                                timeout=timeout, stream=True)
        response.raise_for_status()
        pdf_content = response.content
        # PDFのテキスト抽出
        result = await extractor.extract_text(pdf_content, "pdf")
        if not result or not isinstance(result, list):
            return "PDFからテキストを抽出できませんでした。"
        # ページごとのテキストを結合
        texts = ""
        for page in result:
            _num = page.get("page_number", "?")
            _texts = page.get("texts", "内容を読み取れませんでした")
            texts += f"<page{_num}>\n{_texts}\n</page{_num}>\n"

        if len(texts) > max_length:
            texts = texts[:max_length] + "...(省略されました)"
        return texts if texts else "PDFからテキストを抽出できませんでした。"
    except Timeout:
        return "PDFの読み込みがタイムアウトしました。"
    except RequestException as e:
        return f"PDFの読み込み中にエラーが発生しました: {str(e)}"
    except Exception as e:
        return f"予期せぬエラーが発生しました: {str(e)}"


def html2text(url: str, max_length: int = 20_000, timeout: int = 10) -> str:
    """
    urlを受け取り、htmlを取得、不要なタグを削除し整形した内容をstrで返却する

    Args:
        url: 取得対象のURL
        max_length: 返却するテキストの最大長
        timeout: リクエストのタイムアウト秒数

    Returns:
        整形されたHTML内容、またはエラーメッセージ
    """

    # リクエストヘッダーの設定
    headers = {
        'User-Agent': 'Mozilla/5.0 (compatible; MyWebScraper/1.0)',
        'Accept': 'text/html,application/xhtml+xml,application/xml',
        'Accept-Language': 'ja,en-US;q=0.9,en;q=0.8',
    }

    try:
        # リクエスト実行（タイムアウト、リダイレクト制限あり）
        response = requests.get(
            url,
            headers=headers,
            timeout=timeout,
            allow_redirects=True,
            # max_redirects=5,
            stream=True  # 大きなレスポンスを扱うため
        )

        response.raise_for_status()  # 4xx/5xxエラーの場合は例外を発生

        # エンコーディングの適切な処理
        if response.encoding is None:
            response.encoding = 'utf-8'

        # コンテンツサイズの確認（巨大なレスポンスを防ぐ）
        content_length = response.headers.get('Content-Length')
        if content_length and int(content_length) > 10_000_000:  # 10MBを超える場合
            return "コンテンツが大きすぎます。"

        html_content = response.text

        def clean_html() -> str:
            try:
                soup = BeautifulSoup(html_content, 'html.parser')

                # 不要な要素を削除
                for element in ['script', 'style', 'noscript', 'iframe', 'head', 'meta', 'link']:
                    for tag in soup.find_all(element):
                        tag.decompose()

                # コメントを削除
                for comment in soup.find_all(string=lambda string: isinstance(string, Comment)):
                    comment.extract()

                # divタグを解放
                allowed_tags = {'table', 'tr', 'td', 'th', 'ul', 'ol', 'li',
                                'dl', 'dt', 'dd', 'p', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6'}
                for tag in soup.find_all():
                    if tag.name not in allowed_tags:
                        tag.unwrap()

                # 空の要素を削除（img と br は除く）
                for tag in soup.find_all():
                    if len(tag.get_text(strip=True)) == 0:
                        tag.extract()

                # 全ての属性を削除
                for tag in soup.find_all():
                    tag.attrs = {}

                # 空白行と余分な空白を削除
                cleaned_html = re.sub(
                    r'\n\s*\n', '\n', str(soup), flags=re.MULTILINE)
                cleaned_html = re.sub(r'\s+', ' ', cleaned_html)

                return cleaned_html
            except Exception as e:
                return f"HTMLの処理中にエラーが発生しました: {str(e)}"

        cleaned_content = clean_html()

        # 最大長を超える場合はトリミング
        if len(cleaned_content) > max_length:
            cleaned_content = cleaned_content[:max_length] + "...(省略されました)"

        return cleaned_content

    except Timeout:
        return "URLの読み込みがタイムアウトしました。"
    except TooManyRedirects:
        return "リダイレクトが多すぎます。"
    except RequestException as e:
        return f"URLの読み込み中にエラーが発生しました: {str(e)}"
    except Exception as e:
        return f"予期せぬエラーが発生しました: {str(e)}"


def is_valid_url(url: str) -> bool:
    """URLが有効かつ安全かを検証する"""
    try:
        result = urllib.parse.urlparse(url)
        # スキームがhttpまたはhttpsであることを確認
        return all([result.scheme in ['http', 'https'], result.netloc])
    except ValueError:
        return False


async def confirmResult(messages: list[dict], cse_list: list[dict], system_prompt: str) -> list[dict]:
    """Google/Bingの検索結果の本文からユーザーとのやり取りに関連する情報をAOAIで抽出する関数"""
    new_list = []
    logging.info("calling confirmResult")

    async def process_item(item: dict):
        _messages = deepcopy(messages)
        system_messages = [{"role": "system", "content": system_prompt.format(
            datetime_for_today=date_time_ite,
            snippet=item["snippet"])}] + _messages
        response = await AzureOpenAI(
            messages=system_messages,
            model_name="gpt4.1-mini",
            temperature=0,
            timeout=300,
            registry=LLM_REGISTRY
        )
        item["snippet"] = response["choices"][0]["message"]["content"]
        # snippetの長さが30以上の場合のみ item を返す
        if item["snippet"] and len(item["snippet"]) >= 30:
            return item
        return None

    try:
        logging.info("confirmResult: start processing items")
        # 各itemに対してprocess_itemをタスクとして作成
        tasks = [process_item(item) for item in cse_list]
        logging.info(f"confirmResult: {len(tasks)} tasks created")
        results = await asyncio.gather(*tasks)
        new_list = [item for item in results if item is not None]
    except Exception as e:
        logging.critical(f"Error in confirmResult: {e}")
        return []
    return new_list

########################
# Azure OpenAI utility #
########################


def change_system_content(messages: list, system_content: str):
    """
    system roleのmessageが先頭にある場合は上書き、ない場合は先頭に追加する

    元のmessagesを変更してしまっていた可能性があるため、deepcopy()を追加
    """
    new_messages = deepcopy(messages)
    if new_messages[0]["role"] == "system":
        new_messages[0]["content"] = system_content
    else:
        new_system_message = {"role": "system", "content": system_content}
        new_messages.insert(0, new_system_message)
    return new_messages


def error_response(message):
    response_template = {
        'id': '***********',
        'object': 'error',
        'created': 1234567890,
        'model': 'gpt-4-o',
        'choices': [{
            'index': 0,
            'finish_reason': 'stop',
            'message': {
                'role': 'assistant', 'content': message}}],
        'usage': {'completion_tokens': 0, 'prompt_tokens': 0, 'total_tokens': 0}}
    return response_template

# 仮置き


def change_user_input(messages):
    for i in range(len(messages)//2 - 1):
        if messages[-(i*2 + 3)]["role"] != "user":
            return {"role": "error", "content": "error"}
        user_input = messages[-(i*2 + 3)]["content"]
        if user_input not in ["0", "０", "0:", "０：", "1", "１", "1:", "１："]:
            return messages[-(i*2 + 3)]
    return {"role": "error", "content": "error"}


def check_token(text: str) -> int:
    encoding = tiktoken.get_encoding("cl100k_base")
    token_integers = encoding.encode(text)
    return len(token_integers)


#########
# Brain #
#########
# 仮置き
func_list = [
    {
        "func_name": "LLM_GOOGLE",
        "description": "Google検索を行い、ChatGPTが持っていない情報を取得することができます。為替や株価、人名などを調べることに使えます。ニュースについて尋ねられた場合もこのtoolを選択してください。"
    },
    {
        "func_name": "LLM_CHAT",
        "description": "通常のChatGPTです。要約や翻訳を行う場合はこのtoolを優先的に選択してください。"
    }
]

# system_contentの用意
brain_system_content_template = """Complete the objective as best you can. You have access to the following tools:

{func_list}

# 注意事項
- あなたのタスクは利用するツールを判断するまでです。ユーザーからの質問に回答しないように気をつけてください。
- toolは**必ず**以下のリストの中から選択してください。toolの表記は**必ず**変えないようにしてください。
- ファイルの読み込み、日付の取得は全てのtoolで可能です。

Use the following format:

Task: ユーザーの入力とtoolsからユーザーの望んでいるタスクを推測してください。
Thought: それぞれのtoolのdescriptionを読み、どのtoolを用いることが最適か考えてください。
Tool: 用いるtoolを決定してください。
*************
tools = [{func_names}]
*************
"""

func_names = list(function["func_name"] for function in func_list)
BRAIN_SYSTEM_CONTENT = brain_system_content_template.format(
    func_list=func_list, func_names=func_names)

# usageも後に追加予定？


# list messages, func_list -> str function_name("LLM_CHAT" or "LLM_GOOGLE")
async def Brain(messages, mode="", func_list=func_list):
    # 今までの会話をもとに必要なツールを判断
    # 0:チャット、1:Google
    # messages は 二回分のやりとり5行さかのぼる？
    logging.info('Brain processed a request.')

###########
# 例外処理 #
###########
    # mode指定の場合
    if mode == "chat":
        return "LLM_CHAT"
    elif mode == "google":
        return "LLM_GOOGLE"
    elif mode == "gemini":
        return "LLM_GEMINI"
    elif mode == "inside":
        return "LLM_DOCS"
    elif mode == "minutes":
        return "LLM_MINUTES"
    elif mode == "box":
        return "LLM_BOX"

    # func_listが読み込めなかった場合
    if len(func_list) == 0:
        logging.critical("Brain: No functions")
        return "LLM_CHAT"

    # 0:,1:が入力された場合 # 今後追加するかも
    user_input = next((content["text"] for content in messages[-1]
                      ["content"] if content["type"] == "text"), "")
    if user_input[0] in ("0", "０"):
        return "LLM_CHAT"

    if user_input[0] in ("1", "１"):
        return "LLM_GOOGLE"

    # token数が多い場合の処理
    if check_token(user_input) > 16_000-4096:
        return "LLM_CHAT"

    # model
    model_name = "gpt4.1-mini"

########
# init #
########
    brain_system_content = BRAIN_SYSTEM_CONTENT
    # messagesの整形
    brain_messages = change_system_content(
        messages, system_content=brain_system_content)

########
# main #
########
    try:
        response = await AzureOpenAI(
            messages=brain_messages,
            temperature=0,
            model_name=model_name,
            max_retries=2,
            timeout=230,
            raise_for_error=False,
            registry=LLM_REGISTRY
        )
        response_text = response["choices"][0]["message"]["content"]
    except:
        logging.critical("Brain: main AzureOpenAI ERROR")
        return "LLM_CHAT"
    logging.info("Brain response: " + response_text)

    # BrainModelの出力から選択されたToolを抽出
    pattern = r'Tool:.+'  # 正規表現パターン
    matches = re.findall(pattern, response_text)

    # 正規表現
    try:
        # assert len(matches) > 0
        # マッチの抽出
        func_name = matches[-1][5:]
        # 前後の空白の除去
        func_name = func_name.strip()
        # どのfunctionかを抽出
        func_name = func_name.replace('"', '').replace("'", "")
        assert func_name in func_names

    # 全文照会
    except Exception as e:
        logging.info(f"Brain: Failed Regular Expression, {e}")
        flag = False
        for func_tool in func_list:
            func_name = func_tool["func_name"]
            if func_name in response_text:
                flag = True
                break
        # 一致するものがなかった場合
        if not flag:
            func_name = "LLM_CHAT"
    return func_name


# new function
# ファイルアップロード機能
async def messages2textMessages(upn: str, messages: list, async_http_client: AsyncHttpClient, image_required: bool = True) -> tuple:
    """
    messagesをaoaiへのリクエスト用に整形する
    レスポンス
    - text_messages: 整形後のmessages
    - blobs: アップロードしたファイル名の一覧
    1. messageがテキストのみの場合
        - content = textにする
    2. messageにfileが含まれている場合
        - fileの文字起こしを行う
        - ユニークな名前をつけ、blobに元のファイル、文字起こしファイルをそれぞれアップロードする。
        - contentを整形する
        - 文字起こしファイルをblobsに追加する
    3. messageにblobが指定されている場合
        - blobから指定されたblobをダウンロードする(画像とtxtのみ対応)
        - contentを整形する
    ファイルの文字起こし、blobへのアップロード、messagesの整形
    """
    new_messages = []
    blobs = []

    # Blobサービスクライアントの初期化
    blob_service_client = BlobServiceClient.from_connection_string(
        BLOB_CONNECTION_STRING)
    container_client = blob_service_client.get_container_client(
        container=FILE_CONTAINER_NAME)

    # コンテナが存在しない場合は作成
    if not container_client.exists():
        try:
            container_client.create_container()
        except Exception as e:
            logging.debug(
                f"Container already exists or error creating container: {e}")

    _messages = deepcopy(messages)

    for message in _messages:
        content_list = []  # contentを保持するリスト
        text_list = []  # 文字起こし結果を保持するリスト
        if isinstance(message["content"], str):
            content_list.append({
                "type": "text",
                "text": message["content"]
            })
        elif isinstance(message["content"], list):
            contents = message["content"]
            for content in contents:
                if content["type"] == "text":
                    content_list.append(content)
                elif content["type"] == "file":
                    file_name = content["name"]
                    file_ext = file_name.split(".")[-1].lower()
                    encoded_content = content["data"].split(",")[-1]
                    file_content = base64.b64decode(encoded_content)

                    file_text = ""

                    if file_ext in ("png", "jpg", "jpeg"):
                        # if file_ext in ():
                        if image_required:
                            # AOAIに画像の内容を判断させる
                            image_url = f"data:image/{file_ext};base64,{encoded_content}"

                            # テキストを含まない場合、画像URLをリストに追加
                            content_list.append({
                                "type": "image_url",
                                "image_url": {
                                    "url": image_url
                                }
                            })

                    # 他のファイルタイプの処理
                    content_text, page_num = file2text(file_content, file_ext)
                    file_text = f"## [{file_name}]\n{content_text}"
                    text_list.append(file_text)

                    # ファイルをBlobにアップロード
                    blob_name = upn + "/" + file_name.replace("/", "_")
                    blob_client = blob_service_client.get_blob_client(
                        container=FILE_CONTAINER_NAME, blob=blob_name)

                    base, ext = blob_name.rsplit(".", 1)
                    while blob_client.exists():  # 既存のファイル名の場合は「UUID」を付与
                        blob_name = f"{base}_{uuid.uuid4().hex}.{ext}"
                        blob_client = blob_service_client.get_blob_client(
                            container=FILE_CONTAINER_NAME, blob=blob_name)

                    # ファイルをアップロード
                    blob_client.upload_blob(file_content, overwrite=False)
                    logging.info(f"'{blob_name}' uploaded to blob storage")
                    if image_required:
                        blobs.append(
                            {"name": blob_name, "page": page_num, "file_name": file_name})

                    # txtファイルの作成
                    text_blob_name = ".".join(
                        blob_name.split(".")[:-1]) + ".txt"
                    text_blob_client = blob_service_client.get_blob_client(
                        container=FILE_CONTAINER_NAME, blob=text_blob_name)

                    # ファイルをアップロード # 元のファイルがtxtだった場合用にoverwriteを有効化
                    if not text_blob_client.exists():
                        logging.info(
                            f"'{text_blob_name}' uploaded to blob storage")
                    else:
                        logging.warning(f"'{text_blob_name}' already exists")
                    text_blob_client.upload_blob(file_text, overwrite=True)
                    blobs.append(
                        {"name": text_blob_name, "page": page_num, "file_name": file_name})

                elif content["type"] == "blob":
                    blob_name = content["name"]
                    file_ext = blob_name.split(".")[-1].lower()
                    blob_client = blob_service_client.get_blob_client(
                        container=FILE_CONTAINER_NAME, blob=blob_name)
                    try:
                        if file_ext in ("png", "jpg", "jpeg"):
                            file_content = blob_client.download_blob().readall()
                            content_list.append({
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:image/{file_ext};base64,{base64.b64encode(file_content).decode('utf-8')}"
                                }
                            })
                        elif file_ext == "txt":
                            file_text = blob_client.download_blob().readall().decode('utf-8')
                            text_list.append(file_text)
                        else:
                            continue
                    except Exception as e:
                        logging.warning(f"genie blob: {e}")
                        text_list.append("ファイルの読み込みに失敗しました。")

                elif content["type"] == "link":
                    # boxのリンクを追加
                    content_list.append(content)

        # すべての文字起こし結果を結合
        if text_list:
            content_text = "\n----\n".join(text_list)
            content_list.append({
                "type": "text",
                "text": content_text
            })

        message["content"] = content_list
        new_messages.append(message)
    return new_messages, blobs

# 社内情報とsuffixの除去


def filter_messages(messages, escape_strings="<br><hr>"):
    """
    アシスタントの発言を確認し、文末の定型文の除去を行う。
    文字数が超過していた場合はいままでのやり取りを全削除する。
    """
    tmp_messages = deepcopy(messages)
    new_messages = []
    for message in tmp_messages:
        role = message["role"]
        content = message["content"]

        # assistant以外の場合
        if role != "assistant":
            new_messages.append(message)
            continue

        # 定型文の除去
        new_content = []
        for item in content:
            if item["type"] == "text":
                text = item["text"]
                if escape_strings in text:
                    text_list = text.split(escape_strings)[:-1]
                    text = (escape_strings).join(text_list)
                # 文字数が超過時の処理 # 順番に注意
                if text == "文字数が多すぎます。より短い文でお試しください。":
                    new_messages = []
                    continue
                new_content.append({"type": "text", "text": text})
            elif item["type"] == "image_url":
                # image_urlはそのまま追加
                new_content.append(item)
        message["content"] = new_content
        new_messages.append(message)

    return new_messages


def convert_url_to_a(text):
    """
    テキスト内の全てのURLをHTMLの<a>タグに変換し、target='_blank'を追加する。
    実行順番に注意。
    """
    def convert_html_links():
        """
        <a>タグにtarget='_blank'を追加する。
        既存のtarget属性がある場合は変更を行わない。
        """
        pattern = re.compile(r'(<a\s+(?!.*?target=)[^>]*)(>)')
        return pattern.sub(r"\1 target='_blank'\2", text)

    def convert_markdown_links():
        """
        マークダウン形式のリンク [テキスト](url) をHTMLの<a>タグに変換する。
        """
        pattern = re.compile(r'\[([^\]]+)\]\(([^)]+)\)')
        return pattern.sub(r"<a href='\2' target='_blank'>\1</a>", text)

    def convert_parentheses_links():
        """
        括弧で囲まれたリンク (url) をHTMLの<a>タグに変換する。
        convert_markdown_linksの後に実行。
        """
        pattern = re.compile(r'\((https?://[^\s]+)\)')
        return pattern.sub(r"<a href='\1' target='_blank'>\1</a>", text)

    def convert_plain_urls():
        """
        通常のURLをHTMLの<a>タグに変換する。
        既存の<a>タグがある場合は変更を行わない。
        convert_parentheses_linksの後に実行。
        """
        def process_text_segment(segment):
            """
            通常のURLをHTMLの<a>タグに変換する。
            """
            pattern = re.compile(r'(https?://[^\s]+)')
            return pattern.sub(r"<a href='\1' target='_blank'>\1</a>", segment)

        # テキストを既存の<a>タグで分割し、各セグメントを処理する
        split_segments = re.split(r'(<a[^>]*>.*?</a>)', text, flags=re.DOTALL)
        processed_segments = [
            segment if segment.startswith(
                '<a') else process_text_segment(segment)
            for segment in split_segments
        ]

        return ''.join(processed_segments)

    # すべての変換関数を順に適用
    for text_func in [convert_html_links, convert_markdown_links, convert_parentheses_links, convert_plain_urls]:
        text = text_func()

    return text


def html_replacer(m):
    """
    htmlの可能性があるコードブロックを除去するためのヘルパー関数
    """
    lang = m.group(1)
    content = m.group(2)
    # langが無い場合またはhtmlの場合は```を取り除く
    if lang is None or lang.lower() == "html":
        return content
    else:
        # それ以外の言語は元のブロックをそのまま返す
        return m.group(0)

# ocr


def insert_newline_corrected(text):
    # 正規表現を利用して、特定のピリオドの使用パターンに対しては改行を追加しない
    # 例外のパターンは
    # 1. 数字に挟まれている
    # 2. 大文字の後ろにある
    # 数字に挟まれたピリオド（小数点）を除外
    decimal_pattern = r'(?<=\d)\.(?=\d)'
    # 大文字の後ろにあるピリオド（略語）を除外
    abbreviation_pattern = r'(?<=[A-Z])\.(?!\s|$)'

    # 除外するべきパターンにマッチした部分を一時的に置換
    def temp_replacement(m): return "#_#" if m.group() == '.' else m.group()
    temp_text = re.sub(decimal_pattern, temp_replacement, text)
    temp_text = re.sub(abbreviation_pattern, temp_replacement, temp_text)

    # 「。」と「.」の後に改行を挿入
    # 一時的に置換した部分（###）を元のピリオドに戻す
    final_text = temp_text.replace("。", "。\n").replace(
        ".", ".\n").replace("#_#", ".")

    return final_text

# whisper


def remove_repeated_phrases(text):
    pattern = r'(.*?)\1+'
    while True:
        new_text = re.sub(pattern, r'\1', text)
        if new_text == text:
            break
        text = new_text
    return text


def contains_only_standard_characters(text):
    # 正規表現パターン: ひらがな、カタカナ、漢字、数字、英語のみを許可
    pattern = r'^[ぁ-んァ-ヶー一-龯0-9a-zA-Z\s\!\?\'\-\,\.\%\&\$\¥！？「」、。ー…・％＆＄￥々×〇ゞ]+$'

    # テキストが許可された文字のみを含んでいるかどうかをチェック
    return bool(re.match(pattern, text))


def en_word_counter(text):
    text = text.split(" ")
    num_word = len(text)
    return num_word


def whisper_text_filter(text, language):
    # フィルタの作成
    max_len_text = 125  # ja
    max_num_word = 30  # en

    if not contains_only_standard_characters(text):
        logging.warning(f"Bad Char: {text}")
        text = ""

    if language == "ja":
        if len(text) > max_len_text:
            logging.warning(f"Too mach length: {text}")
            text = ""
    if language == "en":
        num_word = en_word_counter(text)
        if num_word > max_num_word:
            logging.warning(f"Too mach words: {text}")
            text = ""
    return text

##
# mail
##


def send_email(subject: str, content: str, send_to=[], attachment_list=[]) -> bool:
    """
    メールを送信するための関数

    ```python
    attachment_list = [
        {
            "file_name": "minutes.txt",
            "file_content":"議事録.....".encode("utf-8")
        },
        {
            "file_name": "transcribe.txt",
            "file_content":"文字起こし.....".encode("utf-8")
        },
        {
            "file_name": "sample.png",
            "file_content": bytes_data,
            "maintype":"image",
            "subtype": "png"
        },
    ]
    ```
    """
    # mailの設定
    # mail_address = "itc_icolleague@itochu.co.jp"
    send_name = "I-Colleague"
    from_addr = VARIABLE_LIST[ENVIRONMENT_SELECTED]["send_mail_address"]
    # send_from = f"{Header(send_name,'iso-2022-jp')}<tokgv-icolleaguesup@itochu.co.jp>"
    send_from = formataddr((str(Header(send_name, "utf-8")), from_addr))

    # SMTPサーバーの設定 (Outlook.comまたはOffice 365の場合)
    smtp_server = os.environ.get("SMTP_SERVER")
    smtp_port = 25

    # EmailMessageの設定
    msg = EmailMessage()
    msg['From'] = send_from
    # msg['From'] = mail_address
    msg['To'] = ", ".join(send_to)
    msg['Subject'] = subject
    msg["Message-ID"] = make_msgid()

    msg.add_alternative(content, subtype='text')

    # 添付ファイルの処理
    for item in attachment_list:

        # ファイルの形式の取得
        maintype = item.get("maintype", "text")
        subtype = item.get("subtype", "plain")

        # BOMの追加
        if subtype.lower() == "csv":
            item["file_content"] = (
                "\ufeff" + item["file_content"].decode("utf-8")).encode("utf-8")

        msg.add_attachment(
            item["file_content"],
            maintype=maintype,
            subtype=subtype,
            filename=item["file_name"]
        )

    # メール送信の実行
    try:
        with smtplib.SMTP(smtp_server, smtp_port) as server:
            server.ehlo()  # サーバーに挨拶
            server.send_message(
                msg,
                # from_addr=mail_address,
                from_addr=send_from,
                to_addrs=send_to
            )

        logging.debug('メールが送信されました')
        return True
    except Exception as e:
        logging.critical('メールの送信に失敗しました: %s', e)
        return False

##
# blob
##


def upload_blob(file_name: str, file_content: str, container_name: str, overwrite: bool = False) -> bool:
    blob_service_client = BlobServiceClient.from_connection_string(
        conn_str=BLOB_CONNECTION_STRING
    )
    # error handling
    status_flag = False

    # コンテナに接続
    container_client = blob_service_client.get_container_client(
        container=container_name
    )

    try:
        # コンテナが存在しない場合は作成
        try:
            container_client.create_container()
        except Exception as e:
            logging.debug(
                f"Container already exists or error creating container: {e}")

        # Blobクライアントを取得
        blob_client = blob_service_client.get_blob_client(
            container=container_name, blob=file_name)

        # ファイルをアップロード
        blob_client.upload_blob(file_content, overwrite=overwrite)
        logging.info(f"'{file_name}' uploaded to blob storage")
        status_flag = True

    except Exception as e:
        logging.critical(f"Error uploading file: {e}")

    return status_flag


def download_blob(file_name: str, container_name: str) -> bytes:
    blob_service_client = BlobServiceClient.from_connection_string(
        conn_str=BLOB_CONNECTION_STRING
    )
    binary_data = b""
    try:
        # Blobクライアントを取得
        blob_client = blob_service_client.get_blob_client(
            container=container_name, blob=file_name)

        # バイナリデータをダウンロード
        binary_data = blob_client.download_blob().readall()

    except Exception as e:
        logging.critical(f"Error downloading file: {e}")

    return binary_data

##
# Cosmos DB
##


def get_container_proxy(database_name: str, container_name: str):
    """
    Cosmos DBのデータベース、コンテナを指定して操作を行うためのcontainer proxyを取得する
    """
    try:
        client = CosmosClient.from_connection_string(
            conn_str=COSMOS_CONNECTION_STRING)
        database = client.get_database_client(database_name)
        container = database.get_container_client(container_name)
        return container
    except Exception as e:
        logging.critical(f"Error in loading Cosmos DB: {e}")
        raise

###
# Document Intelligence
##


def analyze_documents(file_content: bytes) -> tuple[str, int]:
    def table2md(table) -> str:
        data = table.to_dict()
        # 最大の行数と列数を定義
        max_rows, max_cols = data["row_count"], data["column_count"]

        # リスト（テーブル）を初期化
        table = [["" for _ in range(max_cols)] for __ in range(max_rows)]

        # セルの内容を適切な位置に配置
        for cell in data['cells']:
            # セルの開始行と列を取得
            row, col = cell['row_index'], cell['column_index']
            # セルの内容をテーブルに配置
            table[row][col] = cell['content']
        # table
        df = pd.DataFrame(table[1:], columns=table[0])

        # すべてのカラムがNaN値の行を削除
        df = df.dropna(how='all')
        # すべての行がNaN値の列を削除
        df = df.dropna(axis=1, how='all')

        df = df.map(lambda x: x.replace('\n', '<br>')
                    if isinstance(x, str) else x)
        return df.to_markdown(index=False)

    # page_map = []
    def result2text(result) -> tuple[str, int]:
        ocr_text = ""
        for page_num, page in enumerate(result.pages):
            tables_on_page = [table for table in result.tables if
                              table.bounding_regions[0].page_number == page_num + 1]

            # mark all positions of the table spans in the page
            page_offset = page.spans[0].offset
            page_length = page.spans[0].length
            table_chars = [-1] * page_length
            for table_id, table in enumerate(tables_on_page):
                for span in table.spans:
                    # replace all table spans with "table_id" in table_chars array
                    for i in range(span.length):
                        idx = span.offset - page_offset + i
                        if idx >= 0 and idx < page_length:
                            table_chars[idx] = table_id
            # build page text by replacing charcters in table spans with table html
            page_text = f"### Page {page_num + 1} \n"
            added_tables = set()
            for idx, table_id in enumerate(table_chars):
                if table_id == -1:
                    page_text += result.content[page_offset + idx]
                elif not table_id in added_tables:
                    page_text += "\n" + \
                        table2md(tables_on_page[table_id]) + "\n"
                    added_tables.add(table_id)

            page_text += " "
            # page_map.append((page_num, page_text))
            ocr_text += page_text
        return ocr_text, page_num + 1

    # create your `DocumentAnalysisClient` instance and `AzureKeyCredential` variable
    document_analysis_client = DocumentAnalysisClient(
        endpoint=os.environ.get("DI_ENDPOINT"),
        credential=AzureKeyCredential(os.environ.get("DI_API_KEY"))
    )

    poller = document_analysis_client.begin_analyze_document(
        api_version="2023-07-31",
        model_id="prebuilt-layout",
        document=file_content
    )
    result = poller.result()
    return result2text(result)

# OCR


def file2text(file_content: bytes, file_ext: str) -> tuple[str, int]:
    """
    拡張子ごとに別の処理を呼び出す

    特定の処理がない場合はドキュメントインテリジェンスに投げる
    大きいファイルのアップロード時にエラーが起きるため、pandasは廃止
    """
    try:
        file_ext = file_ext.lower()
        if file_ext in ("pptx", ):
            content_info = pptx2text(file_content)
        elif file_ext in ("xlsx", ):
            content_info = xlsx2text(file_content)
        elif file_ext in ("docx", ):
            content_info = docx2text(file_content)
        elif file_ext in ("txt", "csv",):
            content_info = txt2text(file_content)
        elif file_ext in ("msg", ):
            content_info = msg2text(file_content)
        else:  # if file_ext in ("pdf", "png", "jpg", ):
            content_info = img2text(file_content)
    except Exception as e:
        content_info = "ファイルの内容を読み込めませんでした。", 0
        logging.critical(f"illegal file_ext: {file_ext}, error: {e}")
    return content_info


def img2text(file_content: bytes) -> tuple[str, int]:
    # pdf, png, ...
    return analyze_documents(file_content)


def pptx2text(file_content: bytes) -> tuple[str, int]:
    prs = Presentation(io.BytesIO(file_content))
    text = ""
    # スライドをループ処理
    for slide_number, slide in enumerate(prs.slides):
        text += f"スライド {slide_number + 1}" + "\n"
        # スライド内のすべての形状をループ処理
        for shape in slide.shapes:
            if shape.has_text_frame:
                # テキストを出力
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text += run.text + "\n"
    return text, slide_number + 1


def xlsx2list(file_content: bytes) -> list:
    """
    openpyxlでセルごとに読み込みを実施
    10秒でタイムアウトを設定
    """
    time_limit = 10  # タイムアウトの時間（秒）
    char_limit = 120_000 * 3  # 文字数の上限

    start_time = time.perf_counter()
    char_count = 0
    wb = openpyxl.load_workbook(io.BytesIO(file_content), data_only=True)

    sheet_list = []
    for sheet_name in wb.sheetnames:
        try:
            ws = wb[sheet_name]
            content = []
            current_table = []

            for row in ws.iter_rows(values_only=True):
                row_data = [str(cell).strip()
                            if cell is not None else '' for cell in row]

                # タイムアウトと文字数制限のチェック
                char_count += sum(len(cell) for cell in row_data)
                if time.perf_counter() - start_time > time_limit:
                    raise TimeoutError("タイムアウトが発生しました")
                if char_count > char_limit:
                    raise ValueError("文字数の上限を超えました")

                # 表データの判定（2列以上のデータがある行）
                non_empty_cells = [cell for cell in row_data if cell]
                if len(non_empty_cells) > 1:
                    current_table.append(row_data)
                else:
                    # 表データの終了を検出
                    if current_table:
                        content.append(('table', current_table))
                        current_table = []

                    # 非空のセルを文章として扱う
                    if non_empty_cells:
                        content.append(('text', non_empty_cells[0]))

            # 最後の表データを追加
            if current_table:
                content.append(('table', current_table))

            if not content:
                continue  # コンテンツがない場合は次のシートへ

            sheet_set = (sheet_name, content)
            sheet_list.append(sheet_set)

        except (TimeoutError, ValueError) as e:
            error_sheet_name = "エラーメッセージ"
            error_content = [('text', str(e))]
            sheet_set = (error_sheet_name, error_content)
            sheet_list.append(sheet_set)
            break

    return sheet_list


def xlsx2text(file_content: bytes) -> tuple[str, int]:
    sheet_list = xlsx2list(file_content)

    text = ""
    sheet_num = 0

    for sheet_name, content in sheet_list:
        sheet_num += 1
        sheet_text = f"## {sheet_name}\n\n"

        for content_type, data in content:
            if content_type == 'text':
                sheet_text += f"{data}\n\n"
            elif content_type == 'table':
                # 表データの処理
                cols = list(zip(*data))
                cols = [col for col in cols if any(cell != '' for cell in col)]
                if not cols:
                    continue  # すべての列が空の場合はスキップ
                data = list(zip(*cols))  # 転置して元に戻す

                # Markdown形式に変換
                header = data[0]
                separator = ['-' * max(len(cell), 3) for cell in header]
                header = [cell.replace('\n', '').replace('¥n', '')
                          for cell in header]
                sheet_text += '| ' + ' | '.join(header) + ' |\n'
                sheet_text += '| ' + ' | '.join(separator) + ' |\n'
                for row in data[1:]:
                    row = [cell.replace('\n', '<br>').replace(
                        '¥n', '<br>') for cell in row]
                    sheet_text += '| ' + ' | '.join(row) + ' |\n'
                sheet_text += '\n'

        text += sheet_text

    return text, sheet_num


def docx2text(file_content: bytes) -> tuple[str, int]:
    document = docx.Document(io.BytesIO(file_content))
    full_text = []

    def process_paragraph(paragraph):
        return (paragraph.text or "") + "\n"

    def process_table(table):
        table_text = ""
        for row in table.rows:
            row_text = [
                cell.text.strip() if cell.text else "" for cell in row.cells]
            table_text += " | ".join(row_text) + "\n"
        return table_text + "\n"

    for element in document.element.body:
        if isinstance(element, CT_Tbl):
            table = Table(element, document)
            full_text.append(process_table(table))
        elif element.tag.endswith('p'):
            paragraph = Paragraph(element, document)
            full_text.append(process_paragraph(paragraph))

    return "".join(full_text), 1


def txt2text(file_content: bytes) -> tuple[str, int]:
    # 複数のencodingで読み込みを実施
    for encoding in ENCODINGS:
        try:
            text = file_content.decode(encoding=encoding)
            break
        except UnicodeDecodeError as e:
            continue
    return text, 1


def msg2text(file_content: bytes) -> tuple[str, int]:
    msg = extract_msg.Message(io.BytesIO(file_content))

    subject = msg.subject if msg.subject else "No Subject"
    body = msg.body if msg.body else "No Body"
    sender = msg.sender if msg.sender else "Unknown Sender"
    to = msg.to if msg.to else "Unknown Recipient"
    cc = msg.cc if msg.cc else "No CC"
    date = msg.date if msg.date else "Unknown Date"

    attachments = []
    for attachment in msg.attachments:
        file_name = attachment.name
        if not file_name:
            logging.debug("attachment: None")
            continue
        # msgの場合のみ分岐
        file_ext = file_name.split(".")[-1].lower()
        if file_ext == "msg":
            file_content = attachment.data.__bytes__()
        else:
            file_content = attachment.data

        content_text, _ = file2text(file_content, file_ext)
        file_text = f"#### [{file_name}]\n{content_text}"  # ハッシュを2個追加
        attachments.append(file_text)

    attachments_text = "\n----\n".join(
        attachments) if attachments else "No Attachments"

    text = f"""
### Email Information

**Subject**: {subject}

**Date**: {date}

**Body**: {body}

**Sender**: {sender}

**To**: {to}

**CC**: {cc}

**Attachments**:
{attachments_text}
"""
    return text, 1


# Function to encode a local image into data URL
def encoded_image_to_data_url(image_name, encoded_data):
    # Guess the MIME type of the image based on the file extension
    mime_type, _ = guess_type(image_name)
    if mime_type is None:
        mime_type = 'application/octet-stream'  # Default MIME type if none is found
    # Construct the data URL
    return f"data:{mime_type};base64,{encoded_data}"


def date_time_generator():
    jst = ZoneInfo('Asia/Tokyo')
    while True:
        now = datetime.now(jst)
        formatted_date_time = now.strftime('%m-%d_%Hh')
        yield formatted_date_time


date_time_ite = date_time_generator()

# news


def get_news(k: int = 4) -> list:
    """
    新UIのトップページ用
    blobからエクセルを読み込み、新着ニュースk件を抽出する想定
    現状はハードコーディング

    ```
    [
        {
            "date": "9/20",
            "category": "アップデートのお知らせ",
            "title": "I-Colleagueのリニューアル！",
            "content": "I-ColleagueはX月X日に...",
            "link": "https://example.com"
        },
        {},
        {},
        {}
    ]

    ```
    """

    # blobからexcelファイルをダウンロード

    # 必要なシートを読み込み

    # 整形
    pass

    response_list = []
    sampleCategories = ["アップデートのお知らせ", "イベント情報", "eラーニング", "sample category"]
    sampleData = {
        "date": "9/20",
        "category": "アップデートのお知らせ",
        "title": "I-Colleagueのリニューアル！",
        "content": "I-ColleagueはX月X日に...",
        "link": "https://example.com"
    }
    for _ in range(4):
        sampleData["category"] = random.choice(sampleCategories)
        response_list.append(sampleData.copy())
    return response_list


def generate_info_suffix(user_input: str, file_path: str, sheet_name: str, header: int = 1):

    df = pd.read_excel(file_path, sheet_name=sheet_name, header=header)
    df["keywords"] = df["keywords"].apply(ast.literal_eval)

    info_suffix = ""
    for i in range(len(df)):
        flag = False
        keywords = df["keywords"][i]
        for keyword in keywords:
            if keyword.lower() in user_input.lower():
                flag = True
        if flag:
            info_suffix += "<br><a href='{url}' class='custom-link' target='_blank'>{title}</a>".format(
                url=df["url"][i], title=df["title"][i])
    if len(info_suffix) > 0:
        info_suffix = "【関連サイト】" + info_suffix

    return info_suffix
